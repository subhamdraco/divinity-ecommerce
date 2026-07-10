"""
Generate Divinity Impex Hub - Website Design Presentation (PPT)
For client approval before website development.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand palette ──────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)
GOLD       = RGBColor(0xC9, 0xA2, 0x27)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF5, 0xF7)
MID_GRAY   = RGBColor(0x6B, 0x72, 0x80)
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x2E)
TITAN_BLUE = RGBColor(0x00, 0x6B, 0xB3)
RESHU_ROSE = RGBColor(0xC4, 0x5B, 0x7A)
NOVA_TEAL  = RGBColor(0x0E, 0x7C, 0x7B)
RIZWAN_GOLD = RGBColor(0xB8, 0x86, 0x0B)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color, line_color=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, size=14, bold=False,
                 color=DARK_TEXT, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, size=12, color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  •  {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(4)
    return txBox


def slide_header(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(1.1), NAVY)
    add_text_box(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.5),
                 title, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, Inches(0.6), Inches(0.65), Inches(10), Inches(0.35),
                     subtitle, size=13, color=GOLD)
    add_rect(slide, Inches(0), Inches(1.1), SLIDE_W, Inches(0.04), GOLD)


def mockup_navbar(slide, y=Inches(1.4), brand="DIVINITY IMPEX"):
    add_rect(slide, Inches(0.5), y, Inches(12.33), Inches(0.55), WHITE, MID_GRAY)
    add_text_box(slide, Inches(0.7), y + Inches(0.1), Inches(2.5), Inches(0.35),
                 brand, size=11, bold=True, color=NAVY)
    nav_items = ["Home", "Brands", "About", "Leadership", "Contact"]
    for i, item in enumerate(nav_items):
        add_text_box(slide, Inches(8.5 + i * 0.75), y + Inches(0.12),
                     Inches(0.7), Inches(0.3), item, size=9, color=MID_GRAY)


def mockup_card(slide, left, top, width, height, title, desc, accent):
    add_rect(slide, left, top, width, height, WHITE, MID_GRAY, radius=True)
    add_rect(slide, left, top, width, Inches(0.06), accent)
    add_rect(slide, left + Inches(0.15), top + Inches(0.25),
              width - Inches(0.3), Inches(1.2), LIGHT_GRAY, radius=True)
    add_text_box(slide, left + Inches(0.15), top + Inches(1.55),
                 width - Inches(0.3), Inches(0.35), title, size=13, bold=True, color=NAVY)
    add_text_box(slide, left + Inches(0.15), top + Inches(1.9),
                 width - Inches(0.3), Inches(0.8), desc, size=9, color=MID_GRAY)
    add_rect(slide, left + Inches(0.15), top + height - Inches(0.45),
              Inches(1.0), Inches(0.28), accent, radius=True)
    add_text_box(slide, left + Inches(0.15), top + height - Inches(0.43),
                 Inches(1.0), Inches(0.25), "Explore →", size=8, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════

def slide_cover(prs):
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(3.2), SLIDE_W, Inches(0.06), GOLD)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
                 "DIVINITY IMPEX", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(2.3), Inches(11), Inches(0.6),
                 "Global Brands Business Hub", size=24, color=GOLD, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.6), Inches(11), Inches(0.5),
                 "Website Design Presentation", size=20, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.4),
                 "Concept · Architecture · Wireframes · Design System", size=14,
                 color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.4),
                 "Prepared for Review & Approval  |  July 2026", size=11,
                 color=MID_GRAY, align=PP_ALIGN.CENTER)


def slide_overview(prs):
    slide = blank_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    slide_header(slide, "Project Overview", "Unified digital hub for four business verticals")

    segments = [
        ("TITAN CORE", "Sports Supplement Range", TITAN_BLUE,
         "Cutting-edge nutrition & performance products for the modern health enthusiast."),
        ("RESHU", "FMCG / Personal Care Range", RESHU_ROSE,
         "Elegant personal care solutions crafted for daily luxury and skincare excellence."),
        ("NOVA", "Inner Compass", NOVA_TEAL,
         "Experiential leadership & human development — corporates, schools & institutions."),
        ("RIZWAN ADATIA", "Individual Portfolio", RIZWAN_GOLD,
         "Entrepreneur, philanthropist & motivational speaker — purpose-driven leadership."),
    ]
    for i, (name, tag, color, desc) in enumerate(segments):
        x = Inches(0.6 + i * 3.1)
        mockup_card(slide, x, Inches(1.5), Inches(2.85), Inches(3.2), name, desc, color)
        add_text_box(slide, x, Inches(1.35), Inches(2.85), Inches(0.25),
                     tag, size=8, color=color, align=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0.6), Inches(5.0), Inches(12), Inches(0.4),
                 "Goal: A single premium hub website that showcases all four segments and routes visitors to dedicated sub-sites.",
                 size=13, bold=True, color=NAVY)
    add_bullet_list(slide, Inches(0.6), Inches(5.4), Inches(12), Inches(1.5), [
        "One cohesive Divinity Impex brand experience with segment-specific identity on sub-pages",
        "Click any segment card → navigates to its dedicated route (/titan-core, /reshu, /nova, /rizwan-adatia)",
        "Shared navigation, footer, and design language — unique accent colors per vertical",
        "Mobile-first responsive design for global audience across Africa, Asia, Middle East & beyond",
    ], size=11)


def slide_sitemap(prs):
    slide = blank_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    slide_header(slide, "Site Architecture & Routes", "URL structure and page hierarchy")

    # Hub box
    hub = add_rect(slide, Inches(5.2), Inches(1.6), Inches(2.9), Inches(0.7), NAVY, radius=True)
    add_text_box(slide, Inches(5.2), Inches(1.75), Inches(2.9), Inches(0.4),
                 "divinityimpex.com  (Hub)", size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    routes = [
        ("/titan-core", "Titan Core", TITAN_BLUE, Inches(0.8)),
        ("/reshu", "Reshu", RESHU_ROSE, Inches(3.6)),
        ("/nova", "NOVA Inner Compass", NOVA_TEAL, Inches(6.4)),
        ("/rizwan-adatia", "Rizwan Adatia", RIZWAN_GOLD, Inches(9.2)),
    ]
    for path, name, color, x in routes:
        add_rect(slide, x, Inches(3.0), Inches(2.6), Inches(0.55), color, radius=True)
        add_text_box(slide, x, Inches(3.08), Inches(2.6), Inches(0.35),
                     name, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(3.65), Inches(2.6), Inches(0.3),
                     path, size=9, color=color, align=PP_ALIGN.CENTER)
        # connector line
        add_rect(slide, x + Inches(1.1), Inches(2.3), Inches(0.04), Inches(0.7), MID_GRAY)

    # Sub-routes for NOVA
    add_text_box(slide, Inches(5.8), Inches(4.2), Inches(3), Inches(0.3),
                 "NOVA Sub-routes:", size=10, bold=True, color=NOVA_TEAL)
    nova_subs = ["/nova/corporate", "/nova/schools", "/nova/contact"]
    for i, sub in enumerate(nova_subs):
        add_rect(slide, Inches(5.5 + i * 1.5), Inches(4.55), Inches(1.35), Inches(0.35),
                 LIGHT_GRAY, NOVA_TEAL, radius=True)
        add_text_box(slide, Inches(5.5 + i * 1.5), Inches(4.58), Inches(1.35), Inches(0.3),
                     sub, size=7, color=NOVA_TEAL, align=PP_ALIGN.CENTER)

    # Shared pages
    add_text_box(slide, Inches(0.6), Inches(5.2), Inches(5), Inches(0.3),
                 "Shared Hub Pages:", size=11, bold=True, color=NAVY)
    shared = ["/about", "/global-presence", "/manufacturing", "/leadership", "/contact"]
    for i, p in enumerate(shared):
        col, row = i % 3, i // 3
        add_text_box(slide, Inches(0.6 + col * 2.2), Inches(5.55 + row * 0.35),
                     Inches(2.1), Inches(0.3), p, size=9, color=MID_GRAY)

    add_text_box(slide, Inches(7), Inches(5.2), Inches(5.5), Inches(1.5),
                 "Navigation Flow:\n"
                 "1. Visitor lands on Hub homepage\n"
                 "2. Sees four segment cards with distinct branding\n"
                 "3. Clicks a card → full segment experience on dedicated route\n"
                 "4. Breadcrumb + 'Back to Hub' always available",
                 size=10, color=DARK_TEXT)


def slide_homepage(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    slide_header(slide, "Homepage Wireframe", "divinityimpex.com — The Global Brands Business Hub")

    mockup_navbar(slide)

    # Hero
    add_rect(slide, Inches(0.5), Inches(2.05), Inches(12.33), Inches(2.0), NAVY, radius=True)
    add_text_box(slide, Inches(1), Inches(2.3), Inches(7), Inches(0.6),
                 "Global Manufacturing Partner", size=22, bold=True, color=WHITE)
    add_text_box(slide, Inches(1), Inches(2.95), Inches(7), Inches(0.8),
                 "Crafting high-quality personal care, healthcare & FMCG products\n"
                 "that resonate with consumers worldwide.", size=11, color=LIGHT_GRAY)
    add_rect(slide, Inches(1), Inches(3.55), Inches(1.5), Inches(0.35), GOLD, radius=True)
    add_text_box(slide, Inches(1), Inches(3.57), Inches(1.5), Inches(0.3),
                 "Explore Brands", size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Stats bar
    stats = [("40+", "Years Experience"), ("30+", "Countries"), ("4", "Continents"), ("200+", "Brands")]
    for i, (num, label) in enumerate(stats):
        x = Inches(9 + i * 0.9)
        add_text_box(slide, x, Inches(2.4), Inches(0.8), Inches(0.4),
                     num, size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text_box(slide, x, Inches(2.85), Inches(0.8), Inches(0.3),
                     label, size=7, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Segment cards
    add_text_box(slide, Inches(0.5), Inches(4.2), Inches(5), Inches(0.35),
                 "Our Business Segments", size=14, bold=True, color=NAVY)
    segments = [
        ("TITAN CORE", TITAN_BLUE), ("RESHU", RESHU_ROSE),
        ("NOVA", NOVA_TEAL), ("RIZWAN ADATIA", RIZWAN_GOLD),
    ]
    for i, (name, color) in enumerate(segments):
        x = Inches(0.5 + i * 3.1)
        mockup_card(slide, x, Inches(4.6), Inches(2.85), Inches(2.3), name,
                    "Click to explore →", color)

    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
                 "↓ Scroll: Global Presence Map · R&D Process · Volume Edge · Leadership · Footer",
                 size=9, color=MID_GRAY, align=PP_ALIGN.CENTER)


def slide_titan_core(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    slide_header(slide, "Titan Core — /titan-core", "Sports Supplement Range")

    mockup_navbar(slide, brand="TITAN CORE")
    add_rect(slide, Inches(0.5), Inches(2.05), Inches(12.33), Inches(1.6), TITAN_BLUE, radius=True)
    add_text_box(slide, Inches(1), Inches(2.3), Inches(8), Inches(0.5),
                 "Fuel Your Performance", size=24, bold=True, color=WHITE)
    add_text_box(slide, Inches(1), Inches(2.85), Inches(8), Inches(0.5),
                 "Cutting-edge nutrition & performance products for the modern health enthusiast.",
                 size=12, color=LIGHT_GRAY)

    sections = [
        ("Product Range", "Protein · Pre-workout · Recovery · Vitamins"),
        ("Science-Backed", "R&D driven formulations with international certifications"),
        ("Global Distribution", "Available across Africa, Middle East & Asia Pacific"),
        ("Partner With Us", "Distributor & retail partnership opportunities"),
    ]
    for i, (title, desc) in enumerate(sections):
        x = Inches(0.5 + (i % 2) * 6.2)
        y = Inches(3.9 + (i // 2) * 1.5)
        add_rect(slide, x, y, Inches(5.9), Inches(1.3), LIGHT_GRAY, radius=True)
        add_rect(slide, x, y, Inches(0.08), Inches(1.3), TITAN_BLUE)
        add_text_box(slide, x + Inches(0.25), y + Inches(0.15), Inches(5.4), Inches(0.35),
                     title, size=13, bold=True, color=TITAN_BLUE)
        add_text_box(slide, x + Inches(0.25), y + Inches(0.55), Inches(5.4), Inches(0.6),
                     desc, size=10, color=MID_GRAY)

    add_bullet_list(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5), [
        "Accent: Athletic Blue (#006BB3) · Bold typography · Dynamic imagery · Product grid layout",
    ], size=9, color=MID_GRAY)


def slide_reshu(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    slide_header(slide, "Reshu — /reshu", "FMCG / Personal Care Range")

    mockup_navbar(slide, brand="RESHU")
    add_rect(slide, Inches(0.5), Inches(2.05), Inches(12.33), Inches(1.6),
             RESHU_ROSE, radius=True)
    add_text_box(slide, Inches(1), Inches(2.3), Inches(8), Inches(0.5),
                 "Daily Luxury, Skincare Excellence", size=24, bold=True, color=WHITE)
    add_text_box(slide, Inches(1), Inches(2.85), Inches(8), Inches(0.5),
                 "Elegant personal care solutions crafted for everyday indulgence.",
                 size=12, color=LIGHT_GRAY)

    categories = ["Skincare", "Hair Care", "Body Care", "Fragrances", "Hygiene", "Wellness"]
    for i, cat in enumerate(categories):
        x = Inches(0.5 + (i % 3) * 4.1)
        y = Inches(4.0 + (i // 3) * 1.4)
        add_rect(slide, x, y, Inches(3.8), Inches(1.1), WHITE, RESHU_ROSE, radius=True)
        add_rect(slide, x + Inches(0.15), y + Inches(0.15),
                 Inches(1.0), Inches(0.8), LIGHT_GRAY, radius=True)
        add_text_box(slide, x + Inches(1.3), y + Inches(0.35), Inches(2.3), Inches(0.4),
                     cat, size=12, bold=True, color=RESHU_ROSE)

    add_bullet_list(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5), [
        "Also features sister brands: BONJOUR, ESPECIAL, NAOMI · Soft rose palette · Elegant serif accents · Product catalog grid",
    ], size=9, color=MID_GRAY)


def slide_nova(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    slide_header(slide, "NOVA Inner Compass — /nova", "Leadership & Human Development")

    mockup_navbar(slide, brand="NOVA")
    add_rect(slide, Inches(0.5), Inches(2.05), Inches(12.33), Inches(1.6), NOVA_TEAL, radius=True)
    add_text_box(slide, Inches(1), Inches(2.2), Inches(8), Inches(0.5),
                 "Building Better Leaders. Stronger Teams. Better Outcomes.", size=20, bold=True, color=WHITE)
    add_text_box(slide, Inches(1), Inches(2.75), Inches(8), Inches(0.5),
                 "Experiential leadership programs for corporates, schools & institutions.",
                 size=11, color=LIGHT_GRAY)

    # Two audience tabs
    for i, (tab, path) in enumerate([("Corporate", "/nova/corporate"), ("Schools", "/nova/schools")]):
        x = Inches(0.5 + i * 6.3)
        add_rect(slide, x, Inches(4.0), Inches(5.9), Inches(2.5), LIGHT_GRAY, NOVA_TEAL, radius=True)
        add_rect(slide, x, Inches(4.0), Inches(5.9), Inches(0.45), NOVA_TEAL, radius=True)
        add_text_box(slide, x + Inches(0.2), Inches(4.05), Inches(3), Inches(0.35),
                     tab, size=12, bold=True, color=WHITE)
        content = {
            "Corporate": ["Clarity Shift Program", "Inner Clarity · Outer Impact",
                          "8,500+ Sessions · 13.4M+ Participants", "21 Countries · 16 World Records"],
            "Schools": ["Student Character Development", "Teacher & Parent Impact",
                        "Sponsor & Community Model", "Focus · Discipline · Resilience"],
        }
        add_bullet_list(slide, x + Inches(0.2), Inches(4.6), Inches(5.5), Inches(1.8),
                        content[tab], size=9, color=DARK_TEXT)

    # Framework
    add_text_box(slide, Inches(0.5), Inches(6.65), Inches(12), Inches(0.3),
                 "5A Framework: Aware → Align → Activate → Apply → Amplify  |  "
                 "Contact: info@novaicompass.com  |  www.novaicompass.com",
                 size=9, color=NOVA_TEAL)


def slide_rizwan(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    slide_header(slide, "Rizwan Adatia — /rizwan-adatia", "Entrepreneur · Philanthropist · Speaker")

    mockup_navbar(slide, brand="RIZWAN ADATIA")
    add_rect(slide, Inches(0.5), Inches(2.05), Inches(12.33), Inches(1.6), RIZWAN_GOLD, radius=True)
    add_text_box(slide, Inches(1), Inches(2.2), Inches(8), Inches(0.5),
                 "Visionary Entrepreneur, Philanthropist, and Advocate for Positive Change",
                 size=18, bold=True, color=WHITE)
    add_text_box(slide, Inches(1), Inches(2.75), Inches(8), Inches(0.4),
                 '"Success is measured by the lives we uplift."', size=12, color=LIGHT_GRAY)

    pillars = [
        ("COGEF Group", "Global trade across 10 African countries"),
        ("RAF Global", "Philanthropy — empowering millions since 2015"),
        ("Human for Humans", "Global compassion movement"),
        ("Speaking & Media", "56 awards · Motivational speaker · Biopic 'Rizwan'"),
    ]
    for i, (title, desc) in enumerate(pillars):
        x = Inches(0.5 + (i % 2) * 6.2)
        y = Inches(4.0 + (i // 2) * 1.4)
        add_rect(slide, x, y, Inches(5.9), Inches(1.2), LIGHT_GRAY, RIZWAN_GOLD, radius=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.15), Inches(5.5), Inches(0.35),
                     title, size=12, bold=True, color=RIZWAN_GOLD)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.55), Inches(5.5), Inches(0.5),
                     desc, size=10, color=MID_GRAY)

    add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.3),
                 "Gallery · Awards · Media Features · Get Involved  |  www.rizwanadatia.org  |  info@rizwanadatia.com",
                 size=9, color=RIZWAN_GOLD)


def slide_design_system(prs):
    slide = blank_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    slide_header(slide, "Design System", "Typography · Colors · Components · UI Patterns")

    # Colors
    add_text_box(slide, Inches(0.6), Inches(1.4), Inches(3), Inches(0.35),
                 "Color Palette", size=14, bold=True, color=NAVY)
    colors = [
        ("Hub Navy", NAVY, "#0D1B2A"), ("Hub Gold", GOLD, "#C9A227"),
        ("Titan Blue", TITAN_BLUE, "#006BB3"), ("Reshu Rose", RESHU_ROSE, "#C45B7A"),
        ("NOVA Teal", NOVA_TEAL, "#0E7C7B"), ("Rizwan Gold", RIZWAN_GOLD, "#B8860B"),
    ]
    for i, (name, rgb, hex_code) in enumerate(colors):
        x = Inches(0.6 + (i % 3) * 2.2)
        y = Inches(1.85 + (i // 3) * 0.7)
        add_rect(slide, x, y, Inches(0.5), Inches(0.5), rgb, radius=True)
        add_text_box(slide, x + Inches(0.6), y + Inches(0.05), Inches(1.5), Inches(0.2),
                     name, size=9, bold=True, color=DARK_TEXT)
        add_text_box(slide, x + Inches(0.6), y + Inches(0.28), Inches(1.5), Inches(0.2),
                     hex_code, size=8, color=MID_GRAY)

    # Typography
    add_text_box(slide, Inches(7), Inches(1.4), Inches(5), Inches(0.35),
                 "Typography", size=14, bold=True, color=NAVY)
    add_text_box(slide, Inches(7), Inches(1.85), Inches(5.5), Inches(1.5),
                 "Headings: Playfair Display (serif, premium feel)\n"
                 "Body: Inter / Segoe UI (clean, readable)\n"
                 "Accent: Montserrat (CTAs, labels, nav items)\n\n"
                 "Hierarchy: H1 48px · H2 32px · H3 24px · Body 16px",
                 size=10, color=DARK_TEXT)

    # Components
    add_text_box(slide, Inches(0.6), Inches(3.5), Inches(5), Inches(0.35),
                 "UI Components", size=14, bold=True, color=NAVY)
    components = [
        "Segment Cards — hover lift + accent border + route link",
        "Sticky Navigation — hub logo + segment switcher dropdown",
        "Hero Sections — full-width with gradient overlay per segment",
        "Stats Counters — animated numbers on scroll",
        "Breadcrumbs — Hub > Segment > Sub-page",
        "CTA Buttons — Gold (hub) / Segment accent (sub-pages)",
        "Footer — unified contact, social links, segment quick-links",
        "Global Map — interactive SVG for manufacturing hubs",
    ]
    add_bullet_list(slide, Inches(0.6), Inches(3.9), Inches(5.5), Inches(3.0),
                    components, size=10)

    # Responsive
    add_text_box(slide, Inches(7), Inches(3.5), Inches(5.5), Inches(0.35),
                 "Responsive Breakpoints", size=14, bold=True, color=NAVY)
    add_text_box(slide, Inches(7), Inches(3.9), Inches(5.5), Inches(2.5),
                 "Desktop (1280px+): 4-column segment grid\n"
                 "Tablet (768–1279px): 2-column grid, collapsible nav\n"
                 "Mobile (<768px): Single column, hamburger menu\n"
                 "Touch-friendly 44px tap targets\n"
                 "Lazy-loaded images for performance\n"
                 "Smooth page transitions between routes",
                 size=10, color=DARK_TEXT)


def slide_hub_sections(prs):
    slide = blank_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    slide_header(slide, "Hub Page — Additional Sections", "Content from Divinity Impex Global Brands PDF")

    sections = [
        ("Global Vision", "FMCG manufacturer & developer — personal care, healthcare & FMCG"),
        ("Global Presence", "Interactive map: Americas, Africa, Asia Pacific, Europe, Middle East"),
        ("Manufacturing Hubs", "India · China · Turkey · Africa · UAE — sourcing & R&D"),
        ("Trusted Brands", "RESHU, NAOMI, BONJOUR, ESPECIAL, TITAN CORE, RUVO"),
        ("R&D Process", "Data-driven: Category analysis → Sourcing → QA → Volume Edge"),
        ("Legacy", "40+ years · Cash & Carry pillar · Africa distribution backbone"),
        ("Leadership", "Mr. Rizwan Adatia — Purpose-Driven Success philosophy"),
        ("CTA", "Let's Build the Future of FMCG Together"),
    ]
    for i, (title, desc) in enumerate(sections):
        col, row = i % 4, i // 4
        x = Inches(0.5 + col * 3.15)
        y = Inches(1.5 + row * 2.7)
        add_rect(slide, x, y, Inches(2.95), Inches(2.4), WHITE, MID_GRAY, radius=True)
        add_rect(slide, x, y, Inches(2.95), Inches(0.05), GOLD)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.2), Inches(2.65), Inches(0.35),
                     title, size=11, bold=True, color=NAVY)
        add_text_box(slide, x + Inches(0.15), y + Inches(0.6), Inches(2.65), Inches(1.5),
                     desc, size=9, color=MID_GRAY)


def slide_tech_stack(prs):
    slide = blank_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    slide_header(slide, "Technology & Implementation Plan", "Proposed stack for development phase")

    add_text_box(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(3.5),
                 "Proposed Stack:\n\n"
                 "Framework: Next.js 14 (App Router)\n"
                 "Styling: Tailwind CSS + custom design tokens\n"
                 "Animations: Framer Motion (scroll reveals, transitions)\n"
                 "Routing: File-based (/titan-core, /reshu, /nova, /rizwan-adatia)\n"
                 "Deployment: Vercel or custom hosting\n"
                 "SEO: Meta tags, Open Graph per segment\n"
                 "Performance: Image optimization, lazy loading",
                 size=11, color=DARK_TEXT)

    add_text_box(slide, Inches(7), Inches(1.5), Inches(5.5), Inches(3.5),
                 "Development Phases (post-approval):\n\n"
                 "Phase 1: Hub homepage + navigation + design system\n"
                 "Phase 2: Four segment landing pages with content\n"
                 "Phase 3: NOVA sub-routes (corporate, schools)\n"
                 "Phase 4: Hub inner pages (about, manufacturing, etc.)\n"
                 "Phase 5: Polish — animations, SEO, mobile QA\n"
                 "Phase 6: Deploy + handoff documentation",
                 size=11, color=DARK_TEXT)

    add_rect(slide, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.2), NAVY, radius=True)
    add_text_box(slide, Inches(1), Inches(5.7), Inches(11.5), Inches(0.8),
                 "⚡  This presentation is for DESIGN APPROVAL ONLY.\n"
                 "Once you approve the concept, wireframes, colors, and structure — "
                 "we will proceed to build the live website.",
                 size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_next_steps(prs):
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)
    add_rect(slide, Inches(0), Inches(3.5), SLIDE_W, Inches(0.06), GOLD)
    add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.6),
                 "Next Steps — Awaiting Your Approval", size=32, bold=True,
                 color=WHITE, align=PP_ALIGN.CENTER)

    steps = [
        "Review this design presentation",
        "Share feedback on layout, colors, content structure",
        "Approve or request revisions to any segment page",
        "Provide brand assets (logos, product images, photography)",
        "Upon approval → we build the full responsive website",
    ]
    for i, step in enumerate(steps):
        y = Inches(2.0 + i * 0.85)
        add_rect(slide, Inches(3.5), y, Inches(0.45), Inches(0.45), GOLD, radius=True)
        add_text_box(slide, Inches(3.5), y + Inches(0.05), Inches(0.45), Inches(0.35),
                     str(i + 1), size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_text_box(slide, Inches(4.2), y + Inches(0.05), Inches(6), Inches(0.4),
                     step, size=14, color=WHITE)

    add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                 "DIVINITY IMPEX  |  Global Operations  |  Quality First",
                 size=14, color=GOLD, align=PP_ALIGN.CENTER)


def main():
    prs = new_prs()
    slide_cover(prs)
    slide_overview(prs)
    slide_sitemap(prs)
    slide_homepage(prs)
    slide_hub_sections(prs)
    slide_titan_core(prs)
    slide_reshu(prs)
    slide_nova(prs)
    slide_rizwan(prs)
    slide_design_system(prs)
    slide_tech_stack(prs)
    slide_next_steps(prs)

    out = r"C:\Users\shubghos\Projects\divinity-impex-hub\Divinity_Impex_Website_Design_Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
